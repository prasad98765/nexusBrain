"""
RAG Service using Qdrant for document storage and retrieval
"""
import os
import logging
import hashlib
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
import tiktoken

# Document processing imports
try:
    import PyPDF2
    import pdfplumber
    from docx import Document as DocxDocument
    from pptx import Presentation
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    PyPDF2 = None
    pdfplumber = None
    DocxDocument = None
    Presentation = None

# Qdrant imports
try:
    from qdrant_client import QdrantClient
    from qdrant_client.models import Distance, VectorParams, PointStruct
    QDRANT_AVAILABLE = True
except ImportError:
    QDRANT_AVAILABLE = False
    QdrantClient = None

# Embedding model
try:
    from sentence_transformers import SentenceTransformer
    import numpy as np
    ML_AVAILABLE = True
except ImportError:
    ML_AVAILABLE = False
    SentenceTransformer = None
    np = None

logger = logging.getLogger(__name__)

class RAGService:
    def __init__(
        self,
        qdrant_url: Optional[str] = None,
        qdrant_api_key: Optional[str] = None,
        collection_name: str = "documents",
        embedding_model: str = "all-MiniLM-L6-v2",
        chunk_size: int = 350,  # tokens
        chunk_overlap: int = 50  # tokens
    ):
        self.qdrant_url = qdrant_url or os.getenv(
            "QDRANT_URL",
            "https://b98b3cb8-72f9-437c-97a7-31cee495b1d8.us-east-1-1.aws.cloud.qdrant.io:6333"
        )
        self.qdrant_api_key = qdrant_api_key or os.getenv("QDRANT_API_KEY")
        self.collection_name = collection_name
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.embedding_model_name = embedding_model
        
        # Initialize Qdrant client
        self.qdrant_client = None
        if QDRANT_AVAILABLE and self.qdrant_api_key:
            try:
                self.qdrant_client = QdrantClient(
                    url=self.qdrant_url,
                    api_key=self.qdrant_api_key,
                    timeout=30
                )
                logger.info(f"Qdrant client initialized successfully: {self.qdrant_url}")
                
                # Ensure collection exists
                self._ensure_collection()
            except Exception as e:
                logger.error(f"Failed to initialize Qdrant client: {e}")
                self.qdrant_client = None
        else:
            logger.warning("Qdrant not available or API key missing")
        
        # Initialize embedding model (same as Redis cache)
        self.embedding_model = None
        if ML_AVAILABLE:
            try:
                self.embedding_model = SentenceTransformer(self.embedding_model_name)
                logger.info(f"Embedding model '{self.embedding_model_name}' loaded successfully")
            except Exception as e:
                logger.error(f"Failed to load embedding model: {e}")
        
        # Initialize tokenizer for chunking
        try:
            self.tokenizer = tiktoken.get_encoding("cl100k_base")
        except Exception as e:
            logger.warning(f"Failed to load tiktoken encoder: {e}")
            self.tokenizer = None
    
    def _ensure_collection(self):
        """Create Qdrant collection if it doesn't exist with proper indexes"""
        if not self.qdrant_client:
            return
        
        try:
            from qdrant_client.models import PayloadSchemaType
            
            collections = self.qdrant_client.get_collections().collections
            collection_names = [c.name for c in collections]
            
            if self.collection_name not in collection_names:
                # Create collection with proper vector configuration
                vector_size = 384  # all-MiniLM-L6-v2 produces 384-dim vectors
                
                self.qdrant_client.create_collection(
                    collection_name=self.collection_name,
                    vectors_config=VectorParams(
                        size=vector_size,
                        distance=Distance.COSINE
                    )
                )
                logger.info(f"Created Qdrant collection: {self.collection_name}")
            else:
                logger.info(f"Qdrant collection already exists: {self.collection_name}")
            
            # Create payload index for workspace_id to enable filtering
            try:
                self.qdrant_client.create_payload_index(
                    collection_name=self.collection_name,
                    field_name="workspace_id",
                    field_schema=PayloadSchemaType.KEYWORD
                )
                logger.info("Created workspace_id index")
            except Exception as idx_error:
                # Index might already exist
                logger.debug(f"Workspace_id index note: {idx_error}")
            
            # Create payload index for filename to enable filtering and deletion
            try:
                self.qdrant_client.create_payload_index(
                    collection_name=self.collection_name,
                    field_name="filename",
                    field_schema=PayloadSchemaType.KEYWORD
                )
                logger.info("Created filename index")
            except Exception as idx_error:
                # Index might already exist
                logger.debug(f"Filename index note: {idx_error}")
                
        except Exception as e:
            logger.error(f"Failed to ensure collection: {e}")
    
    def _count_tokens(self, text: str) -> int:
        """Count tokens in text"""
        if self.tokenizer:
            return len(self.tokenizer.encode(text))
        else:
            # Fallback: rough approximation (1 token ≈ 4 chars)
            return len(text) // 4
    
    def _chunk_text(self, text: str) -> List[str]:
        """Split text into chunks of ~chunk_size tokens with overlap"""
        if not text.strip():
            return []
        
        # Split by paragraphs first
        paragraphs = text.split('\n\n')
        chunks = []
        current_chunk = ""
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            para_tokens = self._count_tokens(para)
            current_tokens = self._count_tokens(current_chunk)
            
            # If adding this paragraph exceeds chunk size
            if current_tokens + para_tokens > self.chunk_size and current_chunk:
                chunks.append(current_chunk.strip())
                # Start new chunk with overlap
                words = current_chunk.split()
                overlap_words = words[-self.chunk_overlap:] if len(words) > self.chunk_overlap else words
                current_chunk = ' '.join(overlap_words) + '\n\n' + para
            else:
                current_chunk += ('\n\n' if current_chunk else '') + para
        
        # Add remaining chunk
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def _generate_embedding(self, text: str) -> Optional[List[float]]:
        """Generate embedding vector for text"""
        if not self.embedding_model:
            return None
        
        try:
            # Truncate text for logging
            text_preview = text[:100] + "..." if len(text) > 100 else text
            logger.debug(f"Generating embedding for text: {text_preview}")
            
            embedding = self.embedding_model.encode(text)
            embedding_list = embedding.tolist() if hasattr(embedding, 'tolist') else list(embedding)
            
            logger.debug(f"✅ Generated {len(embedding_list)}-dimensional embedding")
            return embedding_list
        except Exception as e:
            logger.error(f"Failed to generate embedding: {e}")
            return None
    
    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        text = ""
        
        # Try pdfplumber first (better for complex PDFs)
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
            if text.strip():
                return text
        except Exception as e:
            logger.warning(f"pdfplumber failed: {e}, trying PyPDF2")
        
        # Fallback to PyPDF2
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
        except Exception as e:
            logger.error(f"PyPDF2 failed: {e}")
        
        return text
    
    def extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = DocxDocument(file_path)
            text = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            return text
        except Exception as e:
            logger.error(f"Failed to extract text from DOCX: {e}")
            return ""
    
    def extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract text from PPTX: {e}")
            return ""
    
    def extract_text_from_txt(self, file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                logger.error(f"Failed to read TXT file: {e}")
                return ""
        except Exception as e:
            logger.error(f"Failed to read TXT file: {e}")
            return ""
    
    def process_document(
        self,
        file_path: str,
        filename: str,
        workspace_id: str,
        file_type: str
    ) -> Tuple[int, Optional[str]]:
        """
        Process document: extract text, chunk, embed, and store in Qdrant
        Returns: (number_of_chunks, error_message)
        """
        if not self.qdrant_client or not self.embedding_model:
            return 0, "RAG service not properly initialized"
        
        # Extract text based on file type
        text = ""
        file_type = file_type.lower()
        
        if file_type == 'pdf':
            text = self.extract_text_from_pdf(file_path)
        elif file_type == 'docx':
            text = self.extract_text_from_docx(file_path)
        elif file_type == 'pptx':
            text = self.extract_text_from_pptx(file_path)
        elif file_type == 'txt':
            text = self.extract_text_from_txt(file_path)
        else:
            return 0, f"Unsupported file type: {file_type}"
        
        if not text.strip():
            return 0, "No text could be extracted from document"
        
        # Chunk the text
        chunks = self._chunk_text(text)
        if not chunks:
            return 0, "Failed to create text chunks"
        
        # Generate embeddings and store in Qdrant
        points = []
        timestamp = datetime.utcnow().isoformat()
        
        for idx, chunk in enumerate(chunks):
            embedding = self._generate_embedding(chunk)
            if not embedding:
                continue
            
            # Generate unique ID for this chunk
            chunk_id = hashlib.md5(f"{workspace_id}_{filename}_{idx}_{timestamp}".encode()).hexdigest()
            
            point = PointStruct(
                id=chunk_id,
                vector=embedding,
                payload={
                    "workspace_id": workspace_id,
                    "filename": filename,
                    "text": chunk,
                    "chunk_index": idx,
                    "timestamp": timestamp
                }
            )
            points.append(point)
        
        if not points:
            return 0, "Failed to generate embeddings for chunks"
        
        # Upsert to Qdrant
        try:
            self.qdrant_client.upsert(
                collection_name=self.collection_name,
                points=points
            )
            logger.info(f"Stored {len(points)} chunks for document: {filename}")
            return len(points), None
        except Exception as e:
            logger.error(f"Failed to store chunks in Qdrant: {e}")
            return 0, f"Failed to store in database: {str(e)}"
    
    def process_raw_text(
        self,
        text: str,
        workspace_id: str,
        title: str = "Raw Text"
    ) -> Tuple[int, Optional[str]]:
        """Process raw text input (same as document but no file)"""
        if not self.qdrant_client or not self.embedding_model:
            return 0, "RAG service not properly initialized"
        
        if not text.strip():
            return 0, "Empty text provided"
        
        # Chunk the text
        chunks = self._chunk_text(text)
        if not chunks:
            return 0, "Failed to create text chunks"
        
        # Generate embeddings and store
        points = []
        timestamp = datetime.utcnow().isoformat()
        
        for idx, chunk in enumerate(chunks):
            embedding = self._generate_embedding(chunk)
            if not embedding:
                continue
            
            chunk_id = hashlib.md5(f"{workspace_id}_{title}_{idx}_{timestamp}".encode()).hexdigest()
            
            point = PointStruct(
                id=chunk_id,
                vector=embedding,
                payload={
                    "workspace_id": workspace_id,
                    "filename": title,
                    "text": chunk,
                    "chunk_index": idx,
                    "timestamp": timestamp
                }
            )
            points.append(point)
        
        if not points:
            return 0, "Failed to generate embeddings"
        
        try:
            self.qdrant_client.upsert(
                collection_name=self.collection_name,
                points=points
            )
            logger.info(f"Stored {len(points)} chunks for text: {title}")
            return len(points), None
        except Exception as e:
            logger.error(f"Failed to store text in Qdrant: {e}")
            return 0, f"Failed to store in database: {str(e)}"
    
    def retrieve_context(
        self,
        query: str,
        workspace_id: str,
        top_k: int = 5,
        similarity_threshold: float = 0.5
    ) -> List[Dict[str, Any]]:
        """
        Retrieve relevant chunks from Qdrant based on query
        Returns list of relevant chunks with metadata
        """
        if not self.qdrant_client or not self.embedding_model:
            logger.warning("RAG service not available for retrieval")
            return []
        
        logger.info(f"🔍 RAG Search - Query: '{query[:100]}...' | Workspace: {workspace_id}")
        logger.info(f"📊 Search params - top_k: {top_k}, threshold: {similarity_threshold}")
        
        # Generate query embedding
        query_embedding = self._generate_embedding(query)
        if not query_embedding:
            logger.error("❌ Failed to generate query embedding")
            return []
        
        logger.info(f"✅ Query embedding generated (dim: {len(query_embedding)})")
        
        try:
            # First, check total documents in workspace
            count_result = self.qdrant_client.count(
                collection_name=self.collection_name,
                count_filter={
                    "must": [
                        {"key": "workspace_id", "match": {"value": workspace_id}}
                    ]
                }
            )
            logger.info(f"📚 Total chunks in workspace: {count_result.count}")
            
            # Search Qdrant with workspace filter
            results = self.qdrant_client.search(
                collection_name=self.collection_name,
                query_vector=query_embedding,
                query_filter={
                    "must": [
                        {"key": "workspace_id", "match": {"value": workspace_id}}
                    ]
                },
                limit=top_k,
                score_threshold=similarity_threshold
            )
            
            logger.info(f"🎯 Qdrant returned {len(results)} results above threshold {similarity_threshold}")
            
            # Log all results with scores
            for idx, result in enumerate(results):
                logger.info(f"  Result {idx+1}: score={result.score:.4f}, file='{result.payload.get('filename')}', chunk={result.payload.get('chunk_index')}")
                logger.debug(f"    Preview: {result.payload.get('text', '')[:100]}...")
            
            # If no results, try without threshold to see what's available
            if len(results) == 0 and count_result.count > 0:
                logger.warning(f"⚠️ No results above threshold {similarity_threshold}. Trying without threshold...")
                results_no_threshold = self.qdrant_client.search(
                    collection_name=self.collection_name,
                    query_vector=query_embedding,
                    query_filter={
                        "must": [
                            {"key": "workspace_id", "match": {"value": workspace_id}}
                        ]
                    },
                    limit=top_k
                )
                logger.info(f"📈 Best available scores:")
                for idx, result in enumerate(results_no_threshold[:5]):
                    logger.info(f"  {idx+1}. score={result.score:.4f}, file='{result.payload.get('filename')}'")
            
            # Format results
            contexts = []
            for result in results:
                contexts.append({
                    "text": result.payload.get("text", ""),
                    "filename": result.payload.get("filename", ""),
                    "similarity": result.score,
                    "chunk_index": result.payload.get("chunk_index", 0)
                })
            
            logger.info(f"✅ Retrieved {len(contexts)} relevant chunks for query")
            return contexts
            
        except Exception as e:
            logger.error(f"❌ Failed to retrieve context from Qdrant: {e}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            return []
    
    def list_documents(self, workspace_id: str) -> List[Dict[str, Any]]:
        """
        List all documents for a workspace
        Returns: List of documents with metadata
        """
        if not self.qdrant_client:
            return []
        
        try:
            # Scroll through all points for this workspace
            documents = {}
            scroll_result = self.qdrant_client.scroll(
                collection_name=self.collection_name,
                scroll_filter={
                    "must": [
                        {"key": "workspace_id", "match": {"value": workspace_id}}
                    ]
                },
                limit=1000,
                with_payload=True,
                with_vectors=False
            )
            
            points = scroll_result[0]  # First element is the list of points
            
            # Group chunks by filename
            for point in points:
                filename = point.payload.get("filename")
                timestamp = point.payload.get("timestamp")
                
                if filename not in documents:
                    documents[filename] = {
                        "filename": filename,
                        "chunks": 0,
                        "timestamp": timestamp,
                        "workspace_id": workspace_id
                    }
                documents[filename]["chunks"] += 1
            
            # Convert to list and sort by timestamp
            doc_list = list(documents.values())
            doc_list.sort(key=lambda x: x.get("timestamp", ""), reverse=True)
            
            logger.info(f"Found {len(doc_list)} documents for workspace {workspace_id}")
            return doc_list
            
        except Exception as e:
            logger.error(f"Failed to list documents: {e}")
            return []
    
    def delete_document(self, workspace_id: str, filename: str) -> Tuple[bool, Optional[str]]:
        """
        Delete all chunks of a document from Qdrant
        Returns: (success, error_message)
        """
        if not self.qdrant_client:
            return False, "RAG service not available"
        
        try:
            from qdrant_client.models import Filter, FieldCondition, MatchValue
            
            # Create proper filter for deletion
            delete_filter = Filter(
                must=[
                    FieldCondition(
                        key="workspace_id",
                        match=MatchValue(value=workspace_id)
                    ),
                    FieldCondition(
                        key="filename",
                        match=MatchValue(value=filename)
                    )
                ]
            )
            
            # Delete all points matching workspace_id and filename
            result = self.qdrant_client.delete(
                collection_name=self.collection_name,
                points_selector=delete_filter
            )
            
            logger.info(f"Deleted document '{filename}' for workspace {workspace_id}. Operation result: {result}")
            return True, None
            
        except Exception as e:
            error_msg = f"Failed to delete document: {str(e)}"
            logger.error(error_msg)
            return False, error_msg

# Global RAG service instance
rag_service = RAGService()
